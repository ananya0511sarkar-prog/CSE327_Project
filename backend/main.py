import os
import io
import json
from dotenv import load_dotenv
from datetime import datetime, timedelta, timezone

from fastapi import FastAPI, HTTPException, Depends,Query
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from sqlalchemy import Boolean
from typing import Optional, Dict, Any
from sqlalchemy import Column, String, Integer, select,Float,or_
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.ext.asyncio import create_async_engine, AsyncSession, async_sessionmaker
from passlib.context import CryptContext
import jwt

# --- SDK CLIENT IMPORTS ---
from openai import OpenAI
from google import genai
from anthropic import Anthropic


import pdfplumber
from fastapi import UploadFile, File

# NEW: added for multi-format study material support and PDF figure fallback
import fitz          # PyMuPDF — renders PDF pages to images (used as a figure/cid-garbage fallback)
import base64         # NEW: encodes rendered page images for JSON transport
import re             # NEW: used to detect (cid:N) glyph-garbage in pdfplumber output
from docx import Document as DocxDocument   # NEW: python-docx — for .docx text extraction
from pptx import Presentation               # NEW: python-pptx — for .pptx text extraction


from fastapi import WebSocket, WebSocketDisconnect, HTTPException, Depends
from typing import Dict, List
from fastapi import WebSocket, WebSocketDisconnect

# Load environment variables
load_dotenv()

app = FastAPI()

# Middleware for CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], 
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ─── LIVE SDK INITIALIZATION ─────────────────────────────────────────
openai_key = os.getenv("OPENAI_API_KEY")
gemini_key = os.getenv("GEMINI_API_KEY")
anthropic_key = os.getenv("ANTHROPIC_API_KEY")

openai_client = OpenAI(api_key=openai_key) if openai_key else None
gemini_client = genai.Client(api_key=gemini_key) if gemini_key else None
anthropic_client = Anthropic(api_key=anthropic_key) if anthropic_key else None

print("\n" + "🔑 SDK INITIALIZATION STATUS " + "="*30)
print(f"OpenAI Client:  {'✅ Ready' if openai_client else '❌ Missing Key (OPENAI_API_KEY)'}")
print(f"Gemini Client:  {'✅ Ready' if gemini_client else '❌ Missing Key (GEMINI_API_KEY)'}")
print(f"Claude Client:  {'✅ Ready' if anthropic_client else '❌ Missing Key (ANTHROPIC_API_KEY)'}")
print("="*60 + "\n")


# ─── DATABASE SETUP ─────────────────────────────────────────────────
DATABASE_URL = os.getenv("DATABASE_URL")
#engine = create_async_engine(DATABASE_URL)
engine = create_async_engine(DATABASE_URL, pool_pre_ping=True, pool_recycle=300)
AsyncSessionLocal = async_sessionmaker(engine, expire_on_commit=False)
Base = declarative_base()


# ─── AUTHENTICATION SETUP ───────────────────────────────────────────
pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
JWT_SECRET = os.getenv("JWT_SECRET")
JWT_ALGORITHM = "HS256"
JWT_EXPIRE_MINUTES = 60

def hash_password(password: str) -> str:
    return pwd_context.hash(password)

def verify_password(plain_password: str, hashed_password: str) -> bool:
    return pwd_context.verify(plain_password, hashed_password)

def create_access_token(data: dict) -> str:
    to_encode = data.copy()
    expire = datetime.now(timezone.utc) + timedelta(minutes=JWT_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, JWT_SECRET, algorithm=JWT_ALGORITHM)

from sqlalchemy import ARRAY
# ─── DATABASE MODELS ────────────────────────────────────────────────
class UserDB(Base):
    __tablename__ = "users"
    id = Column(Integer, primary_key=True, index=True)
    email = Column(String, unique=True, index=True)
    password = Column(String)
    role = Column(String)
    #for dashboard page
    phone_number = Column(String, nullable=True)
    tech_stack = Column(ARRAY(String), nullable=True, default=list)
    summary = Column(String, nullable=True)

class UserProfileOut(BaseModel):
    id: int
    email: str
    role: str
    phone_number: Optional[str] = None
    tech_stack: Optional[list[str]] = []
    summary: Optional[str] = None
    class Config:
        from_attributes = True

class UserProfileUpdate(BaseModel):
    phone_number: Optional[str] = None
    tech_stack: Optional[list[str]] = None
    summary: Optional[str] = None

from sqlalchemy import ForeignKey, DateTime, func

class ProjectDB(Base):
    __tablename__ = "projects"
    id = Column(Integer, primary_key=True, index=True)
    title = Column(String, nullable=False)
    status = Column(String, default="In Progress")
    details = Column(String)
    ai_engine = Column(String)
    owner_id = Column(Integer, ForeignKey("users.id"), nullable=False)
    created_at = Column(DateTime(timezone=True), server_default=func.now())


#for study page:
class StudyDocumentDB(Base):
    __tablename__ = "study_documents"
    id = Column(Integer, primary_key=True, index=True)
    project_id = Column(Integer, ForeignKey("projects.id"), nullable=False)
    lesson_id = Column(Integer, ForeignKey("study_lessons.id"), nullable=True)  # === LESSONS BACKEND ===
    name = Column(String, nullable=False)
    content = Column(String, nullable=False)
    pages = Column(Integer, default=1)
    # NEW: stores a JSON string list of {"page": n, "image_base64": "..."} for
    # PDF pages that needed the PyMuPDF render fallback (empty/cid-garbage text,
    # or figure-heavy pages pdfplumber can't represent as text).
    page_images = Column(String, nullable=True)
    annotations = Column(String, nullable=True)  # NEW: JSON string {"<page_num>": [stroke, ...]}
    created_at = Column(DateTime(timezone=True), server_default=func.now())


class StudyActionDB(Base):
    __tablename__ = "study_actions"
    id = Column(Integer, primary_key=True, index=True)
    project_id = Column(Integer, ForeignKey("projects.id"), nullable=False)
    document_name = Column(String, nullable=True)
    lesson_id = Column(Integer, ForeignKey("study_lessons.id"), nullable=True)  # === LESSONS BACKEND ===
    action_type = Column(String, nullable=False)   # "summarize" | "explain" | "questions" | "chat"
    engine = Column(String, nullable=False)
    snippet_text = Column(String, nullable=False)
    ai_reply = Column(String, nullable=False)
    created_at = Column(DateTime(timezone=True), server_default=func.now())

async def get_db():
    async with AsyncSessionLocal() as db:
        yield db

from fastapi.security import OAuth2PasswordBearer
from fastapi import status

oauth2_scheme = OAuth2PasswordBearer(tokenUrl="/api/auth/login")

async def get_current_user(
    token: str = Depends(oauth2_scheme),
    db: AsyncSession = Depends(get_db)
) -> UserDB:
    credentials_exception = HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Could not validate credentials",
        headers={"WWW-Authenticate": "Bearer"},
    )
    try:
        payload = jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
        email: str = payload.get("sub")
        if email is None:
            raise credentials_exception
    except jwt.PyJWTError:
        raise credentials_exception

    result = await db.execute(select(UserDB).where(UserDB.email == email))
    user = result.scalar_one_or_none()
    if user is None:
        raise credentials_exception
    return user

# === LESSONS BACKEND: new tables for real (non-hardcoded) study sessions ===
class StudyLessonDB(Base):
    __tablename__ = "study_lessons"
    id = Column(Integer, primary_key=True, index=True)
    project_id = Column(Integer, ForeignKey("projects.id"), nullable=False)
    title = Column(String, nullable=False)
    description = Column(String, nullable=True)
    created_at = Column(DateTime(timezone=True), server_default=func.now())


class StudyLessonEntryDB(Base):
    __tablename__ = "study_lesson_entries"
    id = Column(Integer, primary_key=True, index=True)
    lesson_id = Column(Integer, ForeignKey("study_lessons.id"), nullable=False)
    project_id = Column(Integer, ForeignKey("projects.id"), nullable=False)
    title = Column(String, nullable=False)
    type = Column(String, nullable=False)   # "summary" | "explanation" | "practice_qa" | "manual_section"
    snippet = Column(String, nullable=True)
    content = Column(String, nullable=False)
    image_url = Column(String, nullable=True)
    style = Column(String, nullable=True)   # JSON string: {"isBold":.., "isItalic":.., "fontSize":..}
    created_at = Column(DateTime(timezone=True), server_default=func.now())
# === END LESSONS BACKEND: new tables ===



@app.on_event("startup")
async def on_startup():
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)

#async def get_db():
    #async with AsyncSessionLocal() as db:
        #yield db





# ─── UPDATED REQUEST SCHEMAS TO MATCH NEXT.JS FRONTEND ──────────────
class SignUpRequest(BaseModel):
    email: str
    password: str
    role: str

class LoginRequest(BaseModel):
    email: str
    password: str

class ForgotPasswordRequest(BaseModel):
    email: str

class ResetPasswordRequest(BaseModel):
    email: str
    new_password: str

class GenerateQuestionRequest(BaseModel):
    difficulty: str
    topic: str
    project_name: str
    engine: str

class ChatQuestionMetadata(BaseModel):
    question_number: Optional[int] = None   # now optional — sandbox still sends it, study page doesn't need to
    title: str
    difficulty: str

class ChatRequest(BaseModel):
    engine: str
    message: str
    code_context: str
    question_metadata: Optional[ChatQuestionMetadata] = None

class ProjectContext(BaseModel):
    name: str
    topic: str

class EvaluateRequest(BaseModel):
    engine: str
    code_context: str
    difficulty_context: str
    question_id: int          # NEW — which saved question this evaluation belongs to
    question_number: int
    question_title: str
    project_context: ProjectContext


class ProjectCreate(BaseModel):
    title: str
    details: Optional[str] = None
    ai_engine: str

class ProjectOut(BaseModel):
    id: int
    title: str
    status: str
    details: Optional[str]
    ai_engine: str

    class Config:
        from_attributes = True

#for sudy page:
class ProcessSnippetRequest(BaseModel):
    engine: str
    action: str
    snippet: str
    document_name: str
    project_id: str
    lesson_id: Optional[int] = None   # NEW: so saved snippets can be scoped to a lesson

#annotation
class AnnotationStroke(BaseModel):
    points: list[list[float]]
    color: str = "#ff3b30"
    width: float = 3

# NEW: a draggable text note pinned to a normalized (x,y) position on the page image
class AnnotationNote(BaseModel):
    id: str                 # client-generated uuid, so drags/edits target the right note
    x: float                # 0..1, normalized left position
    y: float                # 0..1, normalized top position
    text: str
    color: str = "#facc15"  # sticky-note background color

class AnnotationsSaveRequest(BaseModel):
    page: int
    strokes: list[AnnotationStroke]
    notes: list[AnnotationNote] = []   # NEW


#to save questions in db
class ProjectQuestionDB(Base):
    __tablename__ = "project_questions"
    id = Column(Integer, primary_key=True, index=True)
    project_id = Column(Integer, ForeignKey("projects.id"), nullable=False)
    difficulty = Column(String, nullable=False)  # "Easy" | "Medium" | "Hard"
    number = Column(Integer, nullable=False)      # position within that difficulty
    title = Column(String, nullable=False)
    description = Column(String, nullable=False)
    hint = Column(String, nullable=False)
    starter_code = Column(String, nullable=False)
    created_at = Column(DateTime(timezone=True), server_default=func.now())

class ProjectQuestionOut(BaseModel):
    id: int
    difficulty: str
    number: int
    title: str
    description: str
    hint: str
    starter_code: str

    class Config:
        from_attributes = True


#to save the answers in the db:
class QuestionEvaluationDB(Base):
    __tablename__ = "question_evaluations"
    id = Column(Integer, primary_key=True, index=True)
    question_id = Column(Integer, ForeignKey("project_questions.id"), nullable=False)
    code_snapshot = Column(String, nullable=False)   # the code that was submitted for THIS attempt
    evaluation_text = Column(String, nullable=False) # the LLM's reply for THIS attempt
    engine = Column(String, nullable=True)            # which AI graded it
    created_at = Column(DateTime(timezone=True), server_default=func.now())

class QuestionEvaluationOut(BaseModel):
    id: int
    question_id: int
    code_snapshot: str
    evaluation_text: str
    engine: Optional[str] = None
    created_at: datetime

    class Config:
        from_attributes = True



# === FIX 2: saved summaries/flashcards persistence ===
# Response shape for GET /study/saved-items — one entry per saved
# summarize/explain/questions result. Used to repopulate the "Saved
# Summaries & Flashcards" tab from the DB every time the page loads.
class SavedStudyItemOut(BaseModel):
    id: int
    document_name: Optional[str] = None
    lesson_id: Optional[int] = None   # NEW
    action_type: str
    engine: str
    snippet_text: str
    ai_reply: str
    created_at: datetime

    class Config:
        from_attributes = True



# === LESSONS BACKEND: request/response schemas ===
class LessonCreate(BaseModel):
    title: str
    description: Optional[str] = None

class LessonOut(BaseModel):
    id: int
    title: str
    description: Optional[str] = None
    created_at: datetime

    class Config:
        from_attributes = True


class LessonEntryCreate(BaseModel):
    title: str
    type: str            # "summary" | "explanation" | "practice_qa" | "manual_section"
    snippet: Optional[str] = None
    content: str
    image_url: Optional[str] = None
    style: Optional[Dict[str, Any]] = None   # {"isBold":.., "isItalic":.., "fontSize":..}

class LessonEntryOut(BaseModel):
    id: int
    lesson_id: int
    title: str
    type: str
    snippet: Optional[str] = None
    content: str
    image_url: Optional[str] = None
    style: Optional[Dict[str, Any]] = None
    created_at: datetime

    class Config:
        from_attributes = True
# === END LESSONS BACKEND: schemas ===

# ─── AUTH ENDPOINTS ─────────────────────────────────────────────────
@app.post("/api/auth/signup")
async def signup(payload: SignUpRequest, db: AsyncSession = Depends(get_db)):
    clean_email = payload.email.strip().lower()
    clean_role = payload.role.strip().lower()
    
    result = await db.execute(select(UserDB).where(UserDB.email == clean_email))
    user_exists = result.scalar_one_or_none()
    if user_exists:
        return {"success": False, "message": "An account with this email already exists."}
    
    # BCRYPT — used here:
    new_user = UserDB(email=clean_email, password=hash_password(payload.password), role=clean_role)
    db.add(new_user)
    await db.commit()
    return {"success": True, "message": "Account created successfully!"}

@app.post("/api/auth/login")
async def login(payload: LoginRequest, db: AsyncSession = Depends(get_db)):
    clean_email = payload.email.strip().lower()
    result = await db.execute(select(UserDB).where(UserDB.email == clean_email))
    user = result.scalar_one_or_none()
    
    # BCRYPT — used here:
    if user and verify_password(payload.password, user.password):
        # JWT — used here, and only here:
        token = create_access_token({"sub": user.email, "role": user.role})
        return {"success": True, "token": token, "user": {"email": user.email, "role": user.role}}
    return {"success": False, "message": "Invalid credentials."}

@app.post("/api/auth/forgot-password")
async def forgot_password(payload: ForgotPasswordRequest, db: AsyncSession = Depends(get_db)):
    clean_email = payload.email.strip().lower()
    result = await db.execute(select(UserDB).where(UserDB.email == clean_email))
    user = result.scalar_one_or_none()
    if not user:
        return {"success": False, "message": "Email not found."}
    return {"success": True, "message": "Reset link sent."}

@app.post("/api/auth/reset-password")
async def reset_password(payload: ResetPasswordRequest, db: AsyncSession = Depends(get_db)):
    clean_email = payload.email.strip().lower()
    result = await db.execute(select(UserDB).where(UserDB.email == clean_email))
    user = result.scalar_one_or_none()
    if user:
        user.password = hash_password(payload.new_password)
        await db.commit()
        return {"success": True, "message": "Password updated."}
    return {"success": False, "message": "Validation failed."}


# ─── HELPERS FOR RUNNING LIVE API CALLS ──────────────────────────────
def call_live_llm(engine_choice: str, system_instructions: str, user_prompt: str) -> str:
    """Helper utility routing instructions to their matching production API clients."""
    if engine_choice == "ChatGPT":
        if not openai_client:
            raise ValueError("OpenAI client was not initialized. Please check OPENAI_API_KEY.")
            
        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_instructions},
                {"role": "user", "content": user_prompt}
            ]
        )
        return response.choices[0].message.content

    elif engine_choice == "Gemini":
        if not gemini_client:
            raise ValueError("Gemini client was not initialized. Please check GEMINI_API_KEY.")
            
        response = gemini_client.models.generate_content(
            model="gemini-3.1-flash-lite",
            contents=f"{system_instructions}\n\nUser Input:\n{user_prompt}"
        )
        return response.text

    elif engine_choice == "Claude":
        if not anthropic_client:
            raise ValueError("Anthropic client was not initialized. Please check ANTHROPIC_API_KEY.")
            
        response = anthropic_client.messages.create(
            model="claude-3-5-sonnet-20241022",
            max_tokens=1000,
            system=system_instructions,
            messages=[{"role": "user", "content": user_prompt}]
        )
        return response.content[0].text
        
    else:
        raise ValueError("Invalid Engine Choice detected.")


@app.get("/api/projects", response_model=list[ProjectOut])
async def get_projects(
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user)
):
    result = await db.execute(select(ProjectDB).where(ProjectDB.owner_id == current_user.id))
    return result.scalars().all()

@app.post("/api/projects", response_model=ProjectOut)
async def create_project(
    payload: ProjectCreate,
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user)
):
    new_project = ProjectDB(
        title=payload.title,
        details=payload.details,
        ai_engine=payload.ai_engine,
        owner_id=current_user.id
    )
    db.add(new_project)
    await db.commit()
    await db.refresh(new_project)
    return new_project


#for dashboard: 
@app.get("/api/users/me", response_model=UserProfileOut)
async def get_my_profile(current_user: UserDB = Depends(get_current_user)):
    return current_user

@app.patch("/api/users/me", response_model=UserProfileOut)
async def update_my_profile(
    payload: UserProfileUpdate,
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user)
):
    if payload.phone_number is not None:
        current_user.phone_number = payload.phone_number
    if payload.tech_stack is not None:
        current_user.tech_stack = payload.tech_stack
    if payload.summary is not None:
        current_user.summary = payload.summary
    await db.commit()
    await db.refresh(current_user)
    return current_user






@app.get("/api/projects/{project_id}/questions", response_model=list[ProjectQuestionOut])
async def get_project_questions(project_id: str, db: AsyncSession = Depends(get_db)):
    result = await db.execute(
        select(ProjectQuestionDB)
        .where(ProjectQuestionDB.project_id == int(project_id))
        .order_by(ProjectQuestionDB.difficulty, ProjectQuestionDB.number)
    )
    return result.scalars().all()



@app.get("/api/questions/{question_id}/evaluations", response_model=list[QuestionEvaluationOut])
async def get_question_evaluations(question_id: int, db: AsyncSession = Depends(get_db)):
    result = await db.execute(
        select(QuestionEvaluationDB)
        .where(QuestionEvaluationDB.question_id == question_id)
        .order_by(QuestionEvaluationDB.created_at)
    )
    return result.scalars().all()

# ─── INTERACTIVE API SANDBOX ENDPOINTS ───────────────────────────────

@app.post("/api/projects/{project_id}/generate-question")
async def handle_question_generation(
    project_id: str,
    payload: GenerateQuestionRequest,
    db: AsyncSession = Depends(get_db)
):
    try:
        system_context = (
            "You are a coding interview generator. You must generate a single tech question "
            "tailored precisely to the user's project context. You must return your response "
            "EXACTLY as a valid JSON object. Do not include raw markdown wrap strings like ```json.\n\n"
            "The output JSON object structure MUST look like this:\n"
            "{\n"
            '  "title": "Question Name Here",\n'
            '  "description": "Clear step-by-step problem details.",\n'
            '  "hint": "Useful architectural design hint pointer.",\n'
            '  "starterCode": "function solve() {\\n\\n}"\n'
            "}"
        )

        prompt = (
            f"Generate a customized coding challenge with difficulty level '{payload.difficulty}' "
            f"for a project named '{payload.project_name}' focused around this core topic field: '{payload.topic}'."
        )

        raw_reply = call_live_llm(payload.engine, system_context, prompt).strip()

        if raw_reply.startswith("```"):
            lines = raw_reply.splitlines()
            if lines[0].startswith("```"):
                lines = lines[1:]
            if lines[-1].startswith("```"):
                lines = lines[:-1]
            raw_reply = "\n".join(lines).strip()

        parsed_question = json.loads(raw_reply)

        # count existing questions at this difficulty to assign the next number
        result = await db.execute(
            select(ProjectQuestionDB).where(
                ProjectQuestionDB.project_id == int(project_id),
                ProjectQuestionDB.difficulty == payload.difficulty
            )
        )
        existing_count = len(result.scalars().all())

        new_question = ProjectQuestionDB(
            project_id=int(project_id),
            difficulty=payload.difficulty,
            number=existing_count + 1,
            title=parsed_question["title"],
            description=parsed_question["description"],
            hint=parsed_question["hint"],
            starter_code=parsed_question.get("starterCode", "// Write your solution here...")
        )
        db.add(new_question)
        await db.commit()
        await db.refresh(new_question)

        return {"success": True, "question": parsed_question, "question_id": new_question.id}

    except Exception as e:
        print(f"Generation Error Trace: {str(e)}")
        return {"success": False, "error": f"Failed to synthesize valid code schema structures: {str(e)}"}


@app.post("/api/projects/{project_id}/chat")
async def handle_workspace_chat(
    project_id: str,
    payload: ChatRequest,
    db: AsyncSession = Depends(get_db)
):
    try:
        meta_context = ""
        if payload.question_metadata:
            if payload.question_metadata.question_number is not None:
                # Coding sandbox — has a real question number
                meta_context = (
                    f"They are working on Question #{payload.question_metadata.question_number}: "
                    f"\"{payload.question_metadata.title}\" under [{payload.question_metadata.difficulty}] constraints.\n"
                )
            else:
                # Study page — just a document reference, no question number
                meta_context = f"{payload.question_metadata.title}\n"

        system_context = (
            "You are an expert technical interview mentor. The user is writing code in their "
            "Answering Sandbox workspace right now.\n"
            f"{meta_context}"
            f"Here is their current live script context:\n```\n{payload.code_context}\n```\n"
            "Answer their chat query briefly, provide guidance rather than doing it all for them, and focus on clean algorithmic code design patterns."
        )

        ai_reply = call_live_llm(payload.engine, system_context, payload.message)

        new_action = StudyActionDB(
            project_id=int(project_id),
            document_name=None,
            action_type="chat",
            engine=payload.engine,
            snippet_text=payload.message,
            ai_reply=ai_reply
        )
        db.add(new_action)
        await db.commit()

        return {"success": True, "reply": ai_reply}

    except Exception as e:
        print(f"API Error Trace: {str(e)}")
        return {"success": False, "reply": f"API Engine Error: {str(e)}"}

@app.post("/api/projects/{project_id}/evaluate")
async def handle_workspace_evaluation(
    project_id: str,
    payload: EvaluateRequest,
    db: AsyncSession = Depends(get_db)
):
    try:
        evaluation_system = (
            "You are an automated code testing and architecture compiler evaluator. "
            f"Review the code context specifically against Question #{payload.question_number}: '{payload.question_title}'.\n"
            f"Target Context Track: {payload.project_context.name} ({payload.project_context.topic})\n"
            f"Difficulty Baseline Evaluation Level: [{payload.difficulty_context}]\n\n"
            "Provide explicit markdown analysis detailing optimization criteria, runtime scale complexity alternatives, "
            "or structural bug concerns directly. Use clear formatting structural lists."
        )
        user_eval_prompt = (
            f"Evaluate this script context for submission metrics:\n\n"
            f"```\n{payload.code_context}\n```"
        )

        evaluation_reply = call_live_llm(payload.engine, evaluation_system, user_eval_prompt)

        # Save this attempt so it's visible when the user comes back to this question later
        new_eval = QuestionEvaluationDB(
            question_id=payload.question_id,
            code_snapshot=payload.code_context,
            evaluation_text=evaluation_reply,
            engine=payload.engine
        )
        db.add(new_eval)
        await db.commit()

        return {"success": True, "reply": evaluation_reply}

    except Exception as e:
        print(f"Evaluation API Error Trace: {str(e)}")
        return {"success": False, "reply": f"Internal Evaluation Compilation Breakout: {str(e)}"}

# Run once in Neon SQL Editor before restarting uvicorn:
# ALTER TABLE users ADD COLUMN phone_number VARCHAR;
# ALTER TABLE users ADD COLUMN tech_stack TEXT[];
# ALTER TABLE users ADD COLUMN summary VARCHAR;


# ─── NEW: HELPERS FOR MULTI-FORMAT STUDY DOCUMENT EXTRACTION ─────────

# Threshold logic for deciding a PDF page's pdfplumber text is "unusable"
# (empty, or dominated by (cid:N) glyph-ID garbage from font-encoding issues —
# common in LaTeX/Beamer-exported PDFs missing a proper ToUnicode map).
CID_GARBAGE_PATTERN = re.compile(r"\(cid:\d+\)")
MIN_USABLE_TEXT_LENGTH = 20  # chars

def _page_text_is_unusable(page_text: str) -> bool:
    """Flags a page's extracted text as unusable if it's near-empty, or if
    (cid:N) garbage makes up a large share of the extracted characters."""
    stripped = page_text.strip()
    if len(stripped) < MIN_USABLE_TEXT_LENGTH:
        return True
    cid_matches = CID_GARBAGE_PATTERN.findall(stripped)
    if cid_matches and (sum(len(m) for m in cid_matches) / len(stripped)) > 0.15:
        return True
    return False


def _extract_pdf(file_bytes: bytes):
    """Extracts text page-by-page with pdfplumber AND renders every page as a
    PNG via PyMuPDF, so the viewer shows pages exactly like the source PDF
    (figures, diagrams, math notation included) while keeping the extracted
    text available underneath each page for highlighting/search.
    CHANGED: previously only rendered pages flagged as unreadable (empty or
    cid:-garbled text). Now renders ALL pages, since even 'clean' text
    extraction (like CNN.pdf) still misses embedded diagrams entirely —
    pdfplumber only ever sees text, never pixels."""
    extracted_text = ""
    page_images = []  # NEW: now populated for every page, not just flagged ones

    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        page_count = len(pdf.pages)
        page_texts = [page.extract_text() or "" for page in pdf.pages]

    fitz_doc = fitz.open(stream=file_bytes, filetype="pdf")
    for i, page_text in enumerate(page_texts):
        try:
            fitz_page = fitz_doc.load_page(i)
            # NOTE: dropped dpi from 150 -> 110 since every page renders now,
            # to keep response payload manageable on larger PDFs.
            pix = fitz_page.get_pixmap(dpi=110)
            img_b64 = base64.b64encode(pix.tobytes("png")).decode("utf-8")
            page_images.append({"page": i + 1, "image_base64": img_b64})
        except Exception as render_err:
            print(f"PyMuPDF render failed on page {i + 1}: {render_err}")

        # marker is now inserted for every page (not just unusable ones),
        # followed by the page's real text — frontend still needs the text
        # for highlighting/search even though the image is now always shown.
        extracted_text += f"\n\n[[PAGE_IMAGE:{i + 1}]]\n\n"
        extracted_text += page_text + "\n\n"
    fitz_doc.close()

    return extracted_text.strip(), page_count, page_images


def _extract_docx(file_bytes: bytes):
    """Extracts text AND embedded images from a .docx file, in document order.
    CHANGED: previously text-only. Now walks each paragraph's runs looking for
    embedded image blips (inline pictures), extracts their raw bytes via the
    run's relationship parts, and inserts the same [[PAGE_IMAGE:N]] marker
    tokens the PDF extractor uses — so the frontend renders them identically,
    no separate docx-specific frontend code needed."""
    from docx.oxml.ns import qn

    doc = DocxDocument(io.BytesIO(file_bytes))
    extracted_text = ""
    page_images = []
    image_index = 0

    for para in doc.paragraphs:
        for run in para.runs:
            # a:blip elements are how Word embeds inline images inside a run
            blips = run._element.findall('.//' + qn('a:blip'))
            for blip in blips:
                rId = blip.get(qn('r:embed'))
                if rId:
                    try:
                        image_part = run.part.related_parts[rId]
                        image_index += 1
                        img_b64 = base64.b64encode(image_part.blob).decode("utf-8")
                        page_images.append({"page": image_index, "image_base64": img_b64})
                        extracted_text += f"\n\n[[PAGE_IMAGE:{image_index}]]\n\n"
                    except Exception as img_err:
                        print(f"DOCX image extraction failed: {img_err}")
        if para.text.strip():
            extracted_text += para.text + "\n"

    # page_count still approximates like before — .docx has no fixed page
    # concept until rendered
    page_count = max(1, (len(extracted_text) // 1000) + 1)
    return extracted_text.strip(), page_count, page_images


def _extract_pptx(file_bytes: bytes):
    """Extracts text AND embedded pictures from every slide of a .pptx file.
    CHANGED: previously text-only, page_images always returned empty. Now
    walks each slide's shapes, and for any PICTURE-type shape, extracts its
    raw image bytes via shape.image.blob and inserts a [[PAGE_IMAGE:N]]
    marker — same marker format the PDF/DOCX extractors use, so the frontend
    renders it identically with zero extra frontend code.
    NOTE: python-pptx has no rendering engine, so this can only recover
    inserted picture shapes — it cannot render a full slide as one image the
    way PyMuPDF does for PDF pages. Charts, SmartArt, and shapes built from
    vector drawing objects (not actual picture files) won't be captured this
    way — only literal embedded images will."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(file_bytes))
    slide_texts = []
    page_images = []
    image_index = 0

    for slide_number, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line.strip():
                        lines.append(line)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_index += 1
                    img_b64 = base64.b64encode(shape.image.blob).decode("utf-8")
                    page_images.append({"page": image_index, "image_base64": img_b64})
                    lines.append(f"[[PAGE_IMAGE:{image_index}]]")
                except Exception as img_err:
                    print(f"PPTX image extraction failed on slide {slide_number}: {img_err}")
        if lines:
            slide_texts.append(f"--- Slide {slide_number} ---\n" + "\n".join(lines))

    extracted_text = "\n\n".join(slide_texts)
    page_count = len(prs.slides)
    return extracted_text.strip(), page_count, page_images


def _extract_txt(file_bytes: bytes):
    """NEW: decodes a plain .txt file as UTF-8, falling back to latin-1."""
    try:
        extracted_text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        extracted_text = file_bytes.decode("latin-1", errors="ignore")
    page_count = max(1, (len(extracted_text) // 1000) + 1)
    return extracted_text.strip(), page_count, []


#for study page:
@app.post("/api/projects/{project_id}/study/process-snippet")
async def handle_process_snippet(
    project_id: str,
    payload: ProcessSnippetRequest,
    db: AsyncSession = Depends(get_db)
):
    try:
        action_prompts = {
            "summarize": "Summarize the following text into 2-3 concise bullet points.",
            "explain": "Explain the following concept in simple, plain language, using an analogy if it helps.",
                # CHANGED: now asks for an answer after each question, clearly separated,
    # so students can self-check without having to ask again.
        "questions": (
        "Generate 2-3 short practice questions based on the following text. "
        "After EACH question, provide the answer on a new line prefixed with 'Answer:'. "
        "Format each question as:\n\n"
        "**Question N: [Title]**\n"
        "[question text]\n\n"
        "Answer: [clear, concise answer]\n"
    )
        }
        system_context = action_prompts.get(payload.action, "Process the following text.")

        ai_reply = call_live_llm(payload.engine, system_context, payload.snippet)

        new_action = StudyActionDB(
            project_id=int(project_id),
            document_name=payload.document_name,
            lesson_id=payload.lesson_id,   # NEW: persist which lesson this snippet belongs to
            action_type=payload.action,
            engine=payload.engine,
            snippet_text=payload.snippet,
            ai_reply=ai_reply
        )
        db.add(new_action)
        await db.commit()
        await db.refresh(new_action)  # === FIX 2: need id/created_at below ===

        # === FIX 2: saved summaries/flashcards persistence ===
        # Return the saved row's id + timestamp so the frontend can push this
        # result straight into the "Saved Summaries & Flashcards" list right
        # away, instead of doing a full refetch of /study/saved-items.
        return {
            "success": True,
            "result": ai_reply,
            "action_id": new_action.id,
            "created_at": new_action.created_at.isoformat()
        }

        #return {"success": True, "result": ai_reply}

    except Exception as e:
        print(f"Study Snippet Error Trace: {str(e)}")
        return {"success": False, "error": str(e)}
    


@app.post("/api/projects/{project_id}/study/upload-document")
async def upload_study_document(
    project_id: str,
    file: UploadFile = File(...),
    db: AsyncSession = Depends(get_db)
):
    try:
        file_bytes = await file.read()
        filename = file.filename or ""
        extension = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        # CHANGED: routes to the correct extractor based on file extension instead
        # of assuming every upload is a PDF. Unsupported types now return a clear
        # error instead of silently mis-parsing (which is what caused the old bug).
        if extension == "pdf":
            extracted_text, page_count, page_images = _extract_pdf(file_bytes)
        elif extension == "docx":
            extracted_text, page_count, page_images = _extract_docx(file_bytes)
        elif extension == "pptx":
            extracted_text, page_count, page_images = _extract_pptx(file_bytes)
        elif extension == "txt":
            extracted_text, page_count, page_images = _extract_txt(file_bytes)
        else:
            return {
                "success": False,
                "error": f"Unsupported file type: .{extension}. Supported: PDF, DOCX, PPTX, TXT."
            }

        if not extracted_text.strip():
            extracted_text = "No extractable text found — this may be a scanned document (OCR support coming with the Expert page)."

        new_doc = StudyDocumentDB(
            project_id=int(project_id),
            #lesson_id = Column(Integer, ForeignKey("study_lessons.id"), nullable=True)  # === LESSONS BACKEND ===
            name=file.filename,
            content=extracted_text.strip(),
            pages=page_count,
            # NEW: page-image fallbacks stored as a JSON string (None if there were none)
            page_images=json.dumps(page_images) if page_images else None
        )
        db.add(new_doc)
        await db.commit()
        await db.refresh(new_doc)

        return {
            "success": True,
            "document": {
                "id": new_doc.id,
                "name": new_doc.name,
                "content": new_doc.content,
                "pages": new_doc.pages,
                # NEW: parsed back into a list for the frontend to render next to the text
                "page_images": json.loads(new_doc.page_images) if new_doc.page_images else []
            }
        }

    except Exception as e:
        print(f"Document Upload Error Trace: {str(e)}")
        return {"success": False, "error": str(e)}
    

#ALTER TABLE study_documents ADD COLUMN page_images VARCHAR;

# === FIX 2: saved summaries/flashcards persistence ===
@app.get("/api/projects/{project_id}/study/saved-items", response_model=list[SavedStudyItemOut])
async def get_saved_study_items(project_id: str, db: AsyncSession = Depends(get_db)):
    """
    Returns every summarize/explain/questions result saved for this project
    (plain chat messages are excluded — those aren't "study notes"), newest
    first. The frontend calls this once when the Study page mounts, so the
    tab always reflects what's actually in the DB and survives a refresh or
    a logout/login — it's no longer just component state that resets.
    """
    result = await db.execute(
        select(StudyActionDB)
        .where(
            StudyActionDB.project_id == int(project_id),
            StudyActionDB.action_type.in_(["summarize", "explain", "questions"])
        )
        .order_by(StudyActionDB.created_at.desc())
    )
    return result.scalars().all()



# === LESSONS BACKEND: lesson + entry endpoints ===
@app.get("/api/projects/{project_id}/study/lessons", response_model=list[LessonOut])
async def get_lessons(project_id: str, db: AsyncSession = Depends(get_db)):
    result = await db.execute(
        select(StudyLessonDB)
        .where(StudyLessonDB.project_id == int(project_id))
        .order_by(StudyLessonDB.created_at)
    )
    return result.scalars().all()


@app.post("/api/projects/{project_id}/study/lessons", response_model=LessonOut)
async def create_lesson(project_id: str, payload: LessonCreate, db: AsyncSession = Depends(get_db)):
    new_lesson = StudyLessonDB(
        project_id=int(project_id),
        title=payload.title,
        description=payload.description
    )
    db.add(new_lesson)
    await db.commit()
    await db.refresh(new_lesson)
    return new_lesson


@app.get("/api/projects/{project_id}/study/lessons/{lesson_id}/entries", response_model=list[LessonEntryOut])
async def get_lesson_entries(project_id: str, lesson_id: int, db: AsyncSession = Depends(get_db)):
    result = await db.execute(
        select(StudyLessonEntryDB)
        .where(
            StudyLessonEntryDB.project_id == int(project_id),
            StudyLessonEntryDB.lesson_id == lesson_id
        )
        .order_by(StudyLessonEntryDB.created_at)
    )
    entries = result.scalars().all()

    # style is stored as a JSON string in the DB; parse it back to a dict for the frontend
    output = []
    for entry in entries:
        output.append(LessonEntryOut(
            id=entry.id,
            lesson_id=entry.lesson_id,
            title=entry.title,
            type=entry.type,
            snippet=entry.snippet,
            content=entry.content,
            image_url=entry.image_url,
            style=json.loads(entry.style) if entry.style else None,
            created_at=entry.created_at
        ))
    return output


@app.post("/api/projects/{project_id}/study/lessons/{lesson_id}/entries", response_model=LessonEntryOut)
async def create_lesson_entry(
    project_id: str,
    lesson_id: int,
    payload: LessonEntryCreate,
    db: AsyncSession = Depends(get_db)
):
    new_entry = StudyLessonEntryDB(
        project_id=int(project_id),
        lesson_id=lesson_id,
        title=payload.title,
        type=payload.type,
        snippet=payload.snippet,
        content=payload.content,
        image_url=payload.image_url,
        style=json.dumps(payload.style) if payload.style else None
    )
    db.add(new_entry)
    await db.commit()
    await db.refresh(new_entry)

    return LessonEntryOut(
        id=new_entry.id,
        lesson_id=new_entry.lesson_id,
        title=new_entry.title,
        type=new_entry.type,
        snippet=new_entry.snippet,
        content=new_entry.content,
        image_url=new_entry.image_url,
        style=json.loads(new_entry.style) if new_entry.style else None,
        created_at=new_entry.created_at
    )
# === END LESSONS BACKEND: endpoints ===


# NEW: fetches this document's saved annotations, keyed by page number.
# Frontend calls this once when a document becomes active, to repopulate
# strokes/notes drawn in a previous session.
@app.get("/api/projects/{project_id}/study/documents/{document_id}/annotations")
async def get_document_annotations(
    project_id: str, document_id: int, db: AsyncSession = Depends(get_db)
):
    result = await db.execute(select(StudyDocumentDB).where(StudyDocumentDB.id == document_id))
    doc = result.scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    # returns { "1": {"strokes": [...], "notes": [...]}, "2": {...} } — empty dict if nothing saved yet
    return json.loads(doc.annotations) if doc.annotations else {}


@app.put("/api/projects/{project_id}/study/documents/{document_id}/annotations")
async def save_document_annotations(
    project_id: str,
    document_id: int,
    payload: AnnotationsSaveRequest,
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(select(StudyDocumentDB).where(StudyDocumentDB.id == document_id))
    doc = result.scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")

    existing = json.loads(doc.annotations) if doc.annotations else {}
    existing[str(payload.page)] = {
        "strokes": [s.dict() for s in payload.strokes],
        "notes": [n.dict() for n in payload.notes],   # NEW
    }
    doc.annotations = json.dumps(existing)
    await db.commit()
    return {"success": True}


# ─── EXPERT DASHBOARD DATABASE MODELS ───────────────────────────────
class ExpertProfileDB(Base):
    __tablename__ = "expert_profiles"
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"), unique=True, nullable=False)
    name = Column(String, nullable=False)
    email = Column(String, nullable=False)
    role = Column(String, nullable=False)
    company = Column(String, nullable=False)
    avatar = Column(String, nullable=True)
    price_bdt = Column(Integer, default=5000)
    bio = Column(String, nullable=True)
    skills = Column(ARRAY(String), nullable=True, default=list)


class ExpertSlotDB(Base):
    __tablename__ = "expert_slots"
    id = Column(Integer, primary_key=True, index=True)
    expert_id = Column(Integer, ForeignKey("users.id"), nullable=False)
    date = Column(String, nullable=False)
    time = Column(String, nullable=False)
    duration = Column(String, default="45 Mins")
    is_booked = Column(Boolean, default=False)


class ExpertBookingDB(Base):
    __tablename__ = "expert_bookings"
    id = Column(Integer, primary_key=True, index=True)
    expert_id = Column(Integer, ForeignKey("users.id"), nullable=False)
    candidate_name = Column(String, nullable=False)
    candidate_email = Column(String, nullable=False)
    target_role = Column(String, nullable=False)
    date = Column(String, nullable=False)
    time = Column(String, nullable=False)
    status = Column(String, default="Upcoming")
    meeting_url = Column(String, nullable=True)
    fee_bdt = Column(Integer, default=0)


class ExpertWalletDB(Base):
    __tablename__ = "expert_wallets"
    id = Column(Integer, primary_key=True, index=True)
    expert_id = Column(Integer, ForeignKey("users.id"), unique=True, nullable=False)
    total_earned_bdt = Column(Integer, default=0)
    available_withdraw_bdt = Column(Integer, default=0)
    pending_bdt = Column(Integer, default=0)


# ─── EXPERT SCHEMAS ─────────────────────────────────────────────────
class ExpertProfileSchema(BaseModel):
    name: str
    email: str
    role: str
    company: str
    avatar: Optional[str] = None
    priceBDT: int
    bio: str
    skills: list[str] = []


class CreateSlotSchema(BaseModel):
    date: str
    time: str
    duration: Optional[str] = "45 Mins"


class PayoutRequestSchema(BaseModel):
    amount: int


# ─── EXPERT & AUTH ENDPOINTS ─────────────────────────────────────────
@app.post("/api/auth/logout")
async def logout():
    return {"success": True, "message": "Logged out successfully"}


@app.get("/api/expert/profile")
async def get_expert_profile(
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user)
):
    result = await db.execute(
        select(ExpertProfileDB).where(ExpertProfileDB.user_id == current_user.id)
    )
    profile = result.scalar_one_or_none()
    if not profile:
        return {}

    return {
        "name": profile.name,
        "email": profile.email,
        "role": profile.role,
        "company": profile.company,
        "avatar": profile.avatar,
        "priceBDT": profile.price_bdt,
        "bio": profile.bio,
        "skills": profile.skills or []
    }


@app.post("/api/expert/profile")
async def create_expert_profile(
    payload: ExpertProfileSchema,
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user)
):
    result = await db.execute(
        select(ExpertProfileDB).where(ExpertProfileDB.user_id == current_user.id)
    )
    profile = result.scalar_one_or_none()

    if profile:
        profile.name = payload.name
        profile.email = payload.email
        profile.role = payload.role
        profile.company = payload.company
        profile.avatar = payload.avatar
        profile.price_bdt = payload.priceBDT
        profile.bio = payload.bio
        profile.skills = payload.skills
    else:
        profile = ExpertProfileDB(
            user_id=current_user.id,
            name=payload.name,
            email=payload.email,
            role=payload.role,
            company=payload.company,
            avatar=payload.avatar,
            price_bdt=payload.priceBDT,
            bio=payload.bio,
            skills=payload.skills
        )
        db.add(profile)

    await db.commit()
    await db.refresh(profile)

    return {
        "name": profile.name,
        "email": profile.email,
        "role": profile.role,
        "company": profile.company,
        "avatar": profile.avatar,
        "priceBDT": profile.price_bdt,
        "bio": profile.bio,
        "skills": profile.skills or []
    }


@app.put("/api/expert/profile")
async def update_expert_profile(
    payload: ExpertProfileSchema,
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user)
):
    result = await db.execute(
        select(ExpertProfileDB).where(ExpertProfileDB.user_id == current_user.id)
    )
    profile = result.scalar_one_or_none()

    if not profile:
        profile = ExpertProfileDB(
            user_id=current_user.id,
            name=payload.name,
            email=payload.email,
            role=payload.role,
            company=payload.company,
            avatar=payload.avatar,
            price_bdt=payload.priceBDT,
            bio=payload.bio,
            skills=payload.skills
        )
        db.add(profile)
    else:
        profile.name = payload.name
        profile.email = payload.email
        profile.role = payload.role
        profile.company = payload.company
        profile.avatar = payload.avatar
        profile.price_bdt = payload.priceBDT
        profile.bio = payload.bio
        profile.skills = payload.skills

    await db.commit()
    await db.refresh(profile)

    return {
        "name": profile.name,
        "email": profile.email,
        "role": profile.role,
        "company": profile.company,
        "avatar": profile.avatar,
        "priceBDT": profile.price_bdt,
        "bio": profile.bio,
        "skills": profile.skills or []
    }


@app.post("/api/expert/upload-avatar")
async def upload_expert_avatar(
    avatar: UploadFile = File(...),
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user)
):
    file_bytes = await avatar.read()
    encoded_img = base64.b64encode(file_bytes).decode("utf-8")
    avatar_url = f"data:{avatar.content_type};base64,{encoded_img}"

    result = await db.execute(
        select(ExpertProfileDB).where(ExpertProfileDB.user_id == current_user.id)
    )
    profile = result.scalar_one_or_none()
    if profile:
        profile.avatar = avatar_url
        await db.commit()

    return {"avatarUrl": avatar_url}


@app.get("/api/expert/slots")
async def get_expert_slots(
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user)
):
    result = await db.execute(
        select(ExpertSlotDB).where(ExpertSlotDB.expert_id == current_user.id)
    )
    slots = result.scalars().all()
    return [
        {
            "id": slot.id,
            "date": slot.date,
            "time": slot.time,
            "duration": slot.duration,
            "isBooked": slot.is_booked
        }
        for slot in slots
    ]


@app.post("/api/expert/slots")
async def add_expert_slot(
    payload: CreateSlotSchema,
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user)
):
    new_slot = ExpertSlotDB(
        expert_id=current_user.id,
        date=payload.date,
        time=payload.time,
        duration=payload.duration or "45 Mins",
        is_booked=False
    )
    db.add(new_slot)
    await db.commit()
    await db.refresh(new_slot)

    return {
        "id": new_slot.id,
        "date": new_slot.date,
        "time": new_slot.time,
        "duration": new_slot.duration,
        "isBooked": new_slot.is_booked
    }


@app.delete("/api/expert/slots/{slot_id}")
async def remove_expert_slot(
    slot_id: int,
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user)
):
    result = await db.execute(
        select(ExpertSlotDB).where(
            ExpertSlotDB.id == slot_id,
            ExpertSlotDB.expert_id == current_user.id
        )
    )
    slot = result.scalar_one_or_none()
    if not slot:
        raise HTTPException(status_code=404, detail="Slot not found")

    await db.delete(slot)
    await db.commit()
    return {"success": True}


@app.get("/api/expert/bookings")
async def get_expert_bookings(
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user)
):
    result = await db.execute(
        select(ExpertBookingDB).where(ExpertBookingDB.expert_id == current_user.id)
    )
    bookings = result.scalars().all()
    return [
        {
            "id": b.id,
            "candidateName": b.candidate_name,
            "candidateEmail": b.candidate_email,
            "targetRole": b.target_role,
            "date": b.date,
            "time": b.time,
            "status": b.status,
            "meetingUrl": b.meeting_url or "",
            "feeBDT": b.fee_bdt
        }
        for b in bookings
    ]


@app.get("/api/expert/wallet")
async def get_expert_wallet(
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user)
):
    result = await db.execute(
        select(ExpertWalletDB).where(ExpertWalletDB.expert_id == current_user.id)
    )
    wallet = result.scalar_one_or_none()
    if not wallet:
        wallet = ExpertWalletDB(
            expert_id=current_user.id,
            total_earned_bdt=0,
            available_withdraw_bdt=0,
            pending_bdt=0
        )
        db.add(wallet)
        await db.commit()
        await db.refresh(wallet)

    return {
        "totalEarnedBDT": wallet.total_earned_bdt,
        "availableWithdrawBDT": wallet.available_withdraw_bdt,
        "pendingBDT": wallet.pending_bdt
    }


@app.post("/api/expert/payout")
async def request_expert_payout(
    payload: PayoutRequestSchema,
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user)
):
    result = await db.execute(
        select(ExpertWalletDB).where(ExpertWalletDB.expert_id == current_user.id)
    )
    wallet = result.scalar_one_or_none()
    if not wallet or wallet.available_withdraw_bdt < payload.amount:
        raise HTTPException(status_code=400, detail="Insufficient funds")

    wallet.available_withdraw_bdt -= payload.amount
    await db.commit()
    return {"success": True}

@app.get("/api/experts")
async def get_all_experts(db: AsyncSession = Depends(get_db)):
    result = await db.execute(select(ExpertProfileDB))
    profiles = result.scalars().all()
    
    return [
        {
            "id": profile.id,
            "name": profile.name,
            "role": profile.role,
            "company": profile.company,
            "avatar": profile.avatar or "https://images.unsplash.com/photo-1534528741775-53994a69daeb?w=150&auto=format&fit=crop&q=80",
            "priceBDT": profile.price_bdt,
            "bio": profile.bio or "",
            "skills": profile.skills or [],
            "rating": 5.0,
            "reviewsCount": 0
        }
        for profile in profiles
    ]
@app.get("/api/experts/{expert_id}/slots")
async def get_public_expert_slots(expert_id: int, db: AsyncSession = Depends(get_db)):
    # 1. Look up ExpertProfile to resolve both profile.id and user_id
    profile_res = await db.execute(
        select(ExpertProfileDB).where(
            or_(ExpertProfileDB.id == expert_id, ExpertProfileDB.user_id == expert_id)
        )
    )
    profile = profile_res.scalar_one_or_none()

    # 2. Build set of matching IDs (e.g. {1, 10})
    target_ids = {expert_id}
    if profile:
        target_ids.add(profile.id)
        target_ids.add(profile.user_id)

    # 3. Query slots matching any of the resolved IDs
    query = select(ExpertSlotDB).where(
        ExpertSlotDB.expert_id.in_(list(target_ids))
    )
    result = await db.execute(query)
    slots = result.scalars().all()

    return [
        {
            "id": slot.id,
            "date": str(slot.date),
            "time": str(slot.time),
            "duration": getattr(slot, "duration", None) or "45 Mins",
            "isBooked": bool(getattr(slot, "is_booked", False))
        }
        for slot in slots
        if not getattr(slot, "is_booked", False)
    ]
    
class BookingDB(Base):
    __tablename__ = "bookings"

    id = Column(Integer, primary_key=True, index=True)
    student_id = Column(Integer, ForeignKey("users.id"), nullable=False)
    expert_id = Column(Integer, ForeignKey("expert_profiles.id"), nullable=False)
    amount = Column(Float, nullable=False)
    currency = Column(String, default="BDT")
    status = Column(String, default="CONFIRMED")
    created_at = Column(DateTime, default=datetime.utcnow)
from pydantic import BaseModel

class PaymentConfirmSchema(BaseModel):
    expert_id: int
    amount: float
    currency: str = "BDT"

@app.post("/api/payment/confirm")
async def confirm_payment(
    data: PaymentConfirmSchema, 
    db: AsyncSession = Depends(get_db),
    current_user: UserDB = Depends(get_current_user) # Auth dependency to get student ID
):
    # 1. Create a new booking record in PostgreSQL
    new_booking = BookingDB(
        student_id=current_user.id,
        expert_id=data.expert_id,
        amount=data.amount,
        currency=data.currency,
        status="CONFIRMED"
    )
    
    # 2. Persist to database
    db.add(new_booking)
    await db.commit()
    await db.refresh(new_booking)
    
    return {
        "message": "Payment confirmed and booking saved successfully",
        "booking_id": new_booking.id,
        "status": new_booking.status
    }

class CreateBookingPayload(BaseModel):
    slot_id: int
    expert_id: Optional[int] = None
    target_role: Optional[str] = "Software Engineer"
    candidate_name: Optional[str] = None
    candidate_email: Optional[str] = None
    fee_bdt: Optional[int] = 5000


@app.post("/api/bookings")
@app.post("/api/experts/{expert_id}/bookings")
@app.post("/api/expert/bookings")
async def create_expert_booking(
    payload: CreateBookingPayload,
    expert_id: Optional[int] = None,
    db: AsyncSession = Depends(get_db),
    current_user: Optional[UserDB] = Depends(get_current_user)
):
    target_id = expert_id or payload.expert_id

    # 1. Fetch requested slot
    slot_res = await db.execute(
        select(ExpertSlotDB).where(ExpertSlotDB.id == payload.slot_id)
    )
    slot = slot_res.scalar_one_or_none()
    if not slot:
        raise HTTPException(status_code=404, detail="Selected slot not found")

    if getattr(slot, "is_booked", False):
        raise HTTPException(status_code=400, detail="This slot has already been booked")

    # 2. Resolve target_id (Profile ID 1 -> User ID 10) for FK constraint
    user_id_for_expert = slot.expert_id
    if target_id:
        profile_res = await db.execute(
            select(ExpertProfileDB).where(
                or_(ExpertProfileDB.id == target_id, ExpertProfileDB.user_id == target_id)
            )
        )
        profile = profile_res.scalar_one_or_none()
        if profile:
            user_id_for_expert = profile.user_id

    # 3. Determine candidate metadata
    cand_name = payload.candidate_name or getattr(current_user, "full_name", None) or (current_user.email if current_user else "Student Candidate")
    cand_email = payload.candidate_email or (current_user.email if current_user else "candidate@example.com")

    # 4. Insert into expert_bookings table
    new_booking = ExpertBookingDB(
        expert_id=user_id_for_expert,  # Points to public.users(id)
        candidate_name=cand_name,
        candidate_email=cand_email,
        target_role=payload.target_role or "Software Engineer",
        date=str(slot.date),
        time=str(slot.time),
        status="Upcoming",
        meeting_url="https://meet.jit.si/interviewx-session",
        fee_bdt=payload.fee_bdt or 5000
    )

    # 5. Mark slot as booked
    slot.is_booked = True

    db.add(new_booking)
    await db.commit()
    await db.refresh(new_booking)

    return {
        "success": True,
        "message": "Booking saved successfully",
        "booking_id": new_booking.id
    }

@app.get("/api/candidate/bookings")
async def get_bookings_by_candidate_email(
    email: str = Query(..., description="Email address to match candidate bookings"),
    db: AsyncSession = Depends(get_db)
):
    # Query database for bookings matching the provided candidate email
    result = await db.execute(
        select(ExpertBookingDB).where(ExpertBookingDB.candidate_email == email)
    )
    bookings = result.scalars().all()

    if not bookings:
        return {
            "success": True,
            "count": 0,
            "message": f"No bookings found for email: {email}",
            "bookings": []
        }

    return {
        "success": True,
        "count": len(bookings),
        "bookings": bookings
    }


#video calling part


# ─── WEBSOCKET ROOM MANAGER FOR WEBRTC SIGNALING ────────────────────────
class ConnectionManager:
    def __init__(self):
        # Dictionary mapping booking_id (room) to a list of active websocket connections
        self.active_rooms: Dict[int, List[WebSocket]] = {}

    async def connect(self, websocket: WebSocket, room_id: int):
        await websocket.accept()
        if room_id not in self.active_rooms:
            self.active_rooms[room_id] = []
        self.active_rooms[room_id].append(websocket)

    def disconnect(self, websocket: WebSocket, room_id: int):
        if room_id in self.active_rooms:
            if websocket in self.active_rooms[room_id]:
                self.active_rooms[room_id].remove(websocket)
            if not self.active_rooms[room_id]:
                del self.active_rooms[room_id]

    async def broadcast_to_others(self, websocket: WebSocket, room_id: int, message: str):
        if room_id in self.active_rooms:
            for connection in self.active_rooms[room_id]:
                if connection != websocket:
                    await connection.send_text(message)

manager = ConnectionManager()

# ─── WEBRTC SIGNALING WEBSOCKET ENDPOINT ────────────────────────────────
@app.websocket("/api/ws/room/{booking_id}")
async def websocket_signaling_endpoint(
    websocket: WebSocket, 
    booking_id: int,
    user_email: str,       # Passed from frontend query params or token
    user_id: int,          # Passed from frontend query params or token
    is_expert: bool        # Boolean flag indicating user role type
):
    # 1. Verify booking existence and ownership using SQLAlchemy AsyncSession
    async with AsyncSessionLocal() as db:
        result = await db.execute(
            select(ExpertBookingDB).where(ExpertBookingDB.id == booking_id)
        )
        booking = result.scalar_one_or_none()
    
    if not booking:
        await websocket.close(code=4004, reason="Booking not found")
        return

    db_expert_id = booking.expert_id
    db_candidate_email = booking.candidate_email

    # 2. Check authorization rules
    authorized = False
    if is_expert and int(user_id) == int(db_expert_id):
        authorized = True
    elif not is_expert and user_email.lower().strip() == db_candidate_email.lower().strip():
        authorized = True

    if not authorized:
        await websocket.close(code=4003, reason="Unauthorized access to this meeting room")
        return

    # 3. Accept and register connection to the room
    await manager.connect(websocket, booking_id)
    try:
        while True:
            # Receive WebRTC offer/answer/ICE candidate signals and broadcast to peer
            data = await websocket.receive_text()
            await manager.broadcast_to_others(websocket, booking_id, data)
    except WebSocketDisconnect:
        manager.disconnect(websocket, booking_id)
        # Notify other peer that user left
        await manager.broadcast_to_others(websocket, booking_id, '{"type": "peer-disconnected"}')